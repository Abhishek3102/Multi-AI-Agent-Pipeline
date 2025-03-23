import pptx

def extract_text_from_pptx(ppt_path):
    ppt = pptx.Presentation(ppt_path)
    return "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
